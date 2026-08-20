#!/usr/bin/env python3
"""
ココロナビ ハッカソン発表スライド生成スクリプト
都知事杯オープンデータハッカソン2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================
# デザイン定数
# ============================================

# 16:9 サイズ
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# カラーパレット
COLOR_BEIGE = RGBColor(0xF5, 0xE6, 0xD3)       # 背景（明るい木目/生成り）
COLOR_BEIGE_LIGHT = RGBColor(0xFA, 0xF3, 0xE8)  # 薄いベージュ
COLOR_DARK_GREEN = RGBColor(0x2D, 0x50, 0x16)   # メインテキスト（深緑）
COLOR_LIGHT_GREEN = RGBColor(0x7C, 0xB3, 0x42)  # アクセント（若草色）
COLOR_PALE_ORANGE = RGBColor(0xFF, 0xB7, 0x4D)  # アクセント（淡い橙）
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
COLOR_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)       # カード背景
COLOR_CARD_BORDER = RGBColor(0x7C, 0xB3, 0x42)   # カード罫線（若草色）

# フォント
FONT_JP = "Meiryo"
FONT_JP_BOLD = "Meiryo"


def set_slide_background(slide, color):
    """スライド背景を単色で設定"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """角丸の長方形を追加（紙の質感シミュレート）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=COLOR_DARK_GREEN, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_JP):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return txBox


def add_paragraph(text_frame, text, font_size=18, font_color=COLOR_DARK_GREEN,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name=FONT_JP):
    """既存のテキストフレックスに段落を追加"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_accent_bar(slide, left, top, width, height, color=COLOR_LIGHT_GREEN):
    """アクセントバーを追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_slide_1_title(prs):
    """スライド1: タイトル"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # サブタイトル
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.8),
                 "都民のメンタルヘルスとカウンセリングを、公的支援につなぐ",
                 font_size=20, font_color=COLOR_LIGHT_GREEN, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # メインタイトル
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.5),
                 "不安や悩みを、ひとりで抱え込ませない",
                 font_size=40, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # アプリ名
    add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(1.0),
                 "ココロナビ",
                 font_size=52, font_color=COLOR_LIGHT_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # チーム名
    add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.6),
                 "Kirari-TokyoNodeB ｜ 都知事杯オープンデータハッカソン2026",
                 font_size=14, font_color=COLOR_LIGHT_GRAY, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_2_problem(prs):
    """スライド2: 課題"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "必要な支援はある。でも、たどり着きにくい",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # カード背景
    card = add_rounded_rect(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(4.5),
                            COLOR_CARD_BG, COLOR_CARD_BORDER)

    # 課題リスト
    problems = [
        "仕事、生活、人間関係など様々な不安と悩み",
        "就労、生活、心の相談先が分散",
        "自分に合う窓口を判断する根拠が見えにくい"
    ]

    y_pos = Inches(2.2)
    for i, problem in enumerate(problems):
        # 箇条書きアイコン（オレンジドット）
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.8), y_pos + Pt(4), Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_PALE_ORANGE
        dot.line.fill.background()

        add_text_box(slide, Inches(2.2), y_pos, Inches(9), Inches(0.6),
                     problem,
                     font_size=22, font_color=COLOR_DARK_GREEN, bold=False)
        y_pos += Inches(0.8)

    # 結論
    add_accent_bar(slide, Inches(1.5), Inches(5.0), Inches(0.15), Inches(0.8), COLOR_DARK_GREEN)
    add_text_box(slide, Inches(1.8), Inches(5.0), Inches(9), Inches(0.8),
                 "結論：情報の不足ではなく「選べないこと」が課題",
                 font_size=24, font_color=COLOR_DARK_GREEN, bold=True)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_3_solution(prs):
    """スライド3: 解決策（機能一覧）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "気持ちの整理から、次の行動までを一つに",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # 機能カード
    features = [
        ("1", "AI面接練習", "職種に応じた質問生成、回答フィードバック"),
        ("2", "気分チェック・AI傾聴", "気分と状況の記録、共感的な応答と次の行動提案"),
        ("3", "状況別・地域別の支援推薦", "東京都のオープンデータから適合度の高い候補を提示"),
        ("4", "精神科・心療内科検索", "厚生労働省「医療情報ネット」から1,311件を検索"),
        ("5", "母親向け面接対策", "女性特有の状況に応じた助言と面接準備"),
    ]

    y_pos = Inches(1.8)
    for num, title, desc in features:
        # 番号バッジ
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.2), y_pos, Inches(0.6), Inches(0.5)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_LIGHT_GREEN
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_JP
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # タイトル
        add_text_box(slide, Inches(2.0), y_pos - Pt(2), Inches(4), Inches(0.4),
                     title,
                     font_size=18, font_color=COLOR_DARK_GREEN, bold=True)

        # 説明
        add_text_box(slide, Inches(2.0), y_pos + Pt(18), Inches(9), Inches(0.4),
                     desc,
                     font_size=14, font_color=COLOR_LIGHT_GRAY, bold=False)

        y_pos += Inches(0.9)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_4_logic(prs):
    """スライド4: 推荐ロジック"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "推薦は、説明できるルールで決める",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # スコアカード
    scores = [
        ("悩みと支援カテゴリが一致", "60点", COLOR_LIGHT_GREEN),
        ("選択地域と施設所在地が一致", "30点", COLOR_PALE_ORANGE),
        ("相談方法の希望と一致", "15点", COLOR_LIGHT_GREEN),
        ("東京都全域を対象", "12点", COLOR_PALE_ORANGE),
    ]

    y_pos = Inches(1.8)
    for condition, score, accent_color in scores:
        # カード背景
        card = add_rounded_rect(slide, Inches(1.2), y_pos, Inches(10.5), Inches(0.9),
                                COLOR_CARD_BG, COLOR_CARD_BORDER)

        # 条件テキスト
        add_text_box(slide, Inches(1.5), y_pos + Pt(8), Inches(7), Inches(0.5),
                     condition,
                     font_size=20, font_color=COLOR_DARK_GREEN, bold=False)

        # スコアバッジ
        score_badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.5), y_pos + Pt(6), Inches(1.8), Inches(0.5)
        )
        score_badge.fill.solid()
        score_badge.fill.fore_color.rgb = accent_color
        score_badge.line.fill.background()
        tf = score_badge.text_frame
        tf.paragraphs[0].text = score
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_JP
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        y_pos += Inches(1.1)

    # 結論
    add_accent_bar(slide, Inches(1.5), Inches(5.8), Inches(0.15), Inches(0.8), COLOR_DARK_GREEN)
    add_text_box(slide, Inches(1.8), Inches(5.8), Inches(9), Inches(0.8),
                 "生成AIに順位を委ねず、同じ入力には同じ結果を返す",
                 font_size=22, font_color=COLOR_DARK_GREEN, bold=True)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_5_transparency(prs):
    """スライド5: オープンデータの透明性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "オープンデータの根拠を、画面上で確認できる",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # 透明性要素
    elements = [
        "推薦理由を常時表示",
        "出典名と原典リンクを表示",
        "データ更新日を表示",
        "医療機関はデータ基準日を表示",
    ]

    y_pos = Inches(1.8)
    for element in elements:
        # チェックアイコン（若草色）
        check = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), y_pos + Pt(2), Inches(0.4), Inches(0.4)
        )
        check.fill.solid()
        check.fill.fore_color.rgb = COLOR_LIGHT_GREEN
        check.line.fill.background()
        tf = check.text_frame
        tf.paragraphs[0].text = "✓"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_JP
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(slide, Inches(2.2), y_pos, Inches(9), Inches(0.5),
                     element,
                     font_size=20, font_color=COLOR_DARK_GREEN, bold=False)

        y_pos += Inches(0.8)

    # 結論
    add_accent_bar(slide, Inches(1.5), Inches(5.2), Inches(0.15), Inches(0.8), COLOR_DARK_GREEN)
    add_text_box(slide, Inches(1.8), Inches(5.2), Inches(9), Inches(0.8),
                 "便利さだけでなく、透明性と検証可能性を設計",
                 font_size=22, font_color=COLOR_DARK_GREEN, bold=True)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_6_datasource(prs):
    """スライド6: データソース"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "公的データを「行動できる情報」に変えた",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # データソースカード
    sources = [
        ("東京都福祉局", "女性相談、精神保健、社会福祉協議会、福祉人材"),
        ("厚生労働省「医療情報ネット」", "東京23区、精神科・心療内科など1,311件"),
    ]

    y_pos = Inches(1.8)
    for org, desc in sources:
        card = add_rounded_rect(slide, Inches(1.2), y_pos, Inches(10.5), Inches(1.2),
                                COLOR_CARD_BG, COLOR_CARD_BORDER)

        add_text_box(slide, Inches(1.5), y_pos + Pt(8), Inches(9), Inches(0.5),
                     org,
                     font_size=20, font_color=COLOR_DARK_GREEN, bold=True)

        add_text_box(slide, Inches(1.5), y_pos + Pt(32), Inches(9), Inches(0.5),
                     desc,
                     font_size=16, font_color=COLOR_LIGHT_GRAY, bold=False)

        y_pos += Inches(1.5)

    # データ基準日
    add_text_box(slide, Inches(1.2), Inches(5.0), Inches(10.5), Inches(0.5),
                 "データ基準日：医療情報ネット 2026年6月1日",
                 font_size=16, font_color=COLOR_PALE_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_7_conclusion(prs):
    """スライド7: 結論"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "ひとりで抱え込まない心の安心を、東京から",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # メッセージカード
    card = add_rounded_rect(slide, Inches(2.0), Inches(1.8), Inches(9), Inches(3.5),
                            COLOR_CARD_BG, COLOR_CARD_BORDER)

    messages = [
        "AIで気持ちを整理する",
        "公的支援を根拠付きで選ぶ",
        "必要なら医療機関へつながる",
    ]

    y_pos = Inches(2.2)
    for msg in messages:
        # ドット
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(2.8), y_pos + Pt(4), Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_LIGHT_GREEN
        dot.line.fill.background()

        add_text_box(slide, Inches(3.2), y_pos, Inches(7), Inches(0.5),
                     msg,
                     font_size=22, font_color=COLOR_DARK_GREEN, bold=False)
        y_pos += Inches(0.7)

    # 締めのメッセージ
    add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
                 "オープンデータを、一覧表ではなく、迷っている人が行動を選べる体験へ",
                 font_size=20, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # URL
    add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
                 "https://restart-support-ai.pages.dev/",
                 font_size=16, font_color=COLOR_LIGHT_GREEN, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # チーム名
    add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10), Inches(0.5),
                 "Kirari-TokyoNodeB",
                 font_size=14, font_color=COLOR_LIGHT_GRAY, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def create_slide_8_team(prs):
    """スライド8: チーム紹介"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BEIGE)

    # 上部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLOR_LIGHT_GREEN)

    # タイトル
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11), Inches(1.0),
                 "多様な力で、社会課題に挑む",
                 font_size=32, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # チーム情報カード
    card = add_rounded_rect(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                            COLOR_CARD_BG, COLOR_CARD_BORDER)

    team_info = [
        ("CWBJで繋がった非エンジニアチーム", COLOR_LIGHT_GREEN),
        ("各自の専門性や強みを持つ多様なメンバー7名", COLOR_PALE_ORANGE),
        ("AI（Gemini）を活用した開発", COLOR_LIGHT_GREEN),
        ("4つのプロダクトを同時に開発する挑戦", COLOR_PALE_ORANGE),
    ]

    y_pos = Inches(2.2)
    for info, accent_color in team_info:
        # ドット
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(2.2), y_pos + Pt(4), Inches(0.25), Inches(0.25)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_color
        dot.line.fill.background()

        add_text_box(slide, Inches(2.7), y_pos, Inches(8), Inches(0.5),
                     info,
                     font_size=20, font_color=COLOR_DARK_GREEN, bold=False)
        y_pos += Inches(0.8)

    # チーム名
    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
                 "Kirari-TokyoNodeB（ｷﾗﾘﾄｰｷｮｰﾉｰﾄﾞﾋﾞｰ）",
                 font_size=18, font_color=COLOR_DARK_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 下部アクセントバー
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), COLOR_PALE_ORANGE)


def main():
    """メイン関数"""
    # プレゼンテーション作成
    prs = Presentation()

    # 16:9 サイズに設定
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 各スライドを作成
    create_slide_1_title(prs)
    create_slide_2_problem(prs)
    create_slide_3_solution(prs)
    create_slide_4_logic(prs)
    create_slide_5_transparency(prs)
    create_slide_6_datasource(prs)
    create_slide_7_conclusion(prs)
    create_slide_8_team(prs)

    # 保存
    output_path = "presentation/kokoronavi_slides.pptx"
    temp_path = "presentation/kokoronavi_slides_new.pptx"
    prs.save(temp_path)
    print(f"スライドを保存しました: {temp_path}")


if __name__ == "__main__":
    main()
